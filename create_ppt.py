"""
╔══════════════════════════════════════════════════════════════╗
║   Credit Card Fraud Detection — PowerPoint Presentation     ║
║   Generates a professional, university-ready PPT file       ║
╚══════════════════════════════════════════════════════════════╝
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ──────────────────────────────────────────────
#  Color Palette
# ──────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x0C, 0x29)   # Deep navy background
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)   # Card background
ACCENT       = RGBColor(0x6C, 0x5C, 0xE7)   # Purple accent
ACCENT_LIGHT = RGBColor(0xA2, 0x9B, 0xFE)   # Light purple
GREEN        = RGBColor(0x00, 0xD4, 0xAA)   # Teal green
RED          = RGBColor(0xFF, 0x4C, 0x6A)   # Coral red
YELLOW       = RGBColor(0xFD, 0xCB, 0x6E)   # Gold
WHITE        = RGBColor(0xEA, 0xEA, 0xEA)   # Text white
GRAY         = RGBColor(0x88, 0x92, 0xB0)   # Subtitle gray
DARK_TEXT     = RGBColor(0xCC, 0xD6, 0xF6)   # Light blue-white


def set_slide_bg(slide, color=BG_DARK):
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_box(slide, left, top, width, height, fill_color=BG_CARD, border_color=ACCENT):
    """Add a rounded rectangle card to the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bullet_color=ACCENT_LIGHT):
    """Add a bulleted list to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_accent_line(slide, left, top, width):
    """Add a horizontal accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ══════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────
    #  SLIDE 1: Title Slide
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Shield icon
    add_text_box(slide, Inches(4.5), Inches(1.2), Inches(4.5), Inches(1.2),
                 "🛡️", font_size=60, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10.5), Inches(1.2),
                 "Credit Card Fraud Detection", font_size=44,
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(2), Inches(3.5), Inches(9.5), Inches(0.8),
                 "Using Machine Learning", font_size=28,
                 color=ACCENT_LIGHT, bold=False, alignment=PP_ALIGN.CENTER)

    # Accent line
    add_accent_line(slide, Inches(5), Inches(4.4), Inches(3.5))

    # Details
    add_text_box(slide, Inches(2), Inches(4.8), Inches(9.5), Inches(0.5),
                 "Data Science Semester Project — 4th Semester", font_size=18,
                 color=GRAY, alignment=PP_ALIGN.CENTER)

    # Tech badges
    add_text_box(slide, Inches(2), Inches(5.8), Inches(9.5), Inches(0.5),
                 "Python  •  Scikit-learn  •  XGBoost  •  SMOTE  •  Streamlit",
                 font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────
    #  SLIDE 2: Table of Contents
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.8),
                 "📋 Table of Contents", font_size=36, color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(4))

    toc_items = [
        "1.  Introduction & Problem Statement",
        "2.  Dataset Overview",
        "3.  Project Architecture & Tech Stack",
        "4.  Data Preprocessing Pipeline",
        "5.  Exploratory Data Analysis (EDA)",
        "6.  Machine Learning Models",
        "7.  Model Evaluation & Comparison",
        "8.  Results & Best Model Selection",
        "9.  Streamlit Web Dashboard",
        "10. Conclusion & Future Work",
    ]

    for i, item in enumerate(toc_items):
        y = Inches(1.6) + Inches(i * 0.52)
        card = add_shape_box(slide, Inches(1.2), y, Inches(10.5), Inches(0.46))
        add_text_box(slide, Inches(1.5), y + Pt(4), Inches(10), Inches(0.4),
                     item, font_size=17, color=DARK_TEXT)

    # ──────────────────────────────────────────
    #  SLIDE 3: Introduction & Problem Statement
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "1. Introduction & Problem Statement", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    # Problem card
    card = add_shape_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.2))
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "⚠️ The Problem", font_size=22, color=RED, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4),
                    [
                        "• Credit card fraud costs $28+ billion annually worldwide",
                        "• Traditional rule-based systems can't adapt to new fraud patterns",
                        "• Manual review is slow and expensive",
                        "• Fraudulent transactions are extremely rare (< 0.2%)",
                        "• Both false positives and false negatives are costly",
                    ], font_size=15)

    # Solution card
    card = add_shape_box(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2))
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.5),
                 "✅ Our Solution", font_size=22, color=GREEN, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4),
                    [
                        "• Machine Learning for automated fraud detection",
                        "• SMOTE to handle extreme class imbalance",
                        "• Compare 4 ML algorithms to find the best",
                        "• Evaluate using F1 Score (not accuracy!)",
                        "• Deploy via interactive Streamlit dashboard",
                        "• Real-time prediction with confidence scores",
                    ], font_size=15)

    # ──────────────────────────────────────────
    #  SLIDE 4: Dataset Overview
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "2. Dataset Overview", font_size=34, color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Kaggle Credit Card Fraud Detection Dataset — European cardholders, September 2013",
                 font_size=16, color=GRAY)

    # Stats cards
    stats = [
        ("284,807", "Total\nTransactions", WHITE),
        ("492", "Fraudulent\nTransactions", RED),
        ("0.172%", "Fraud\nRate", YELLOW),
        ("31", "Total\nFeatures", GREEN),
        ("2 Days", "Time\nPeriod", ACCENT_LIGHT),
    ]

    for i, (value, label, color) in enumerate(stats):
        x = Inches(0.8 + i * 2.45)
        card = add_shape_box(slide, x, Inches(2.1), Inches(2.2), Inches(1.8))
        add_text_box(slide, x, Inches(2.2), Inches(2.2), Inches(0.9),
                     value, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.1), Inches(2.2), Inches(0.7),
                     label, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Feature description card
    card = add_shape_box(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8))
    add_text_box(slide, Inches(1.1), Inches(4.4), Inches(5), Inches(0.5),
                 "📊 Feature Description", font_size=20, color=ACCENT_LIGHT, bold=True)

    features_left = [
        "• V1 – V28:  PCA-transformed anonymized features",
        "• Time:         Seconds since first transaction",
        "• Amount:     Transaction amount (in Euros)",
        "• Class:        Target → 0 (Legitimate) / 1 (Fraud)",
    ]
    add_bullet_list(slide, Inches(1.1), Inches(5.0), Inches(5.5), Inches(2),
                    features_left, font_size=15)

    key_notes = [
        "• Original features hidden for confidentiality (PCA applied)",
        "• Only Time & Amount retain original meaning",
        "• V1–V28 are already scaled (no extra scaling needed)",
        "• Extreme imbalance: 99.83% legitimate vs 0.17% fraud",
    ]
    add_bullet_list(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(2),
                    key_notes, font_size=15)

    # ──────────────────────────────────────────
    #  SLIDE 5: Project Architecture
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "3. Project Architecture & Tech Stack", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5.5))

    # Architecture flow
    steps = [
        ("📥", "Raw Data\n(CSV)", GRAY),
        ("→", "", ACCENT),
        ("🧹", "Preprocessing\n& Cleaning", GREEN),
        ("→", "", ACCENT),
        ("📊", "EDA\n(5 Plots)", YELLOW),
        ("→", "", ACCENT),
        ("🤖", "Train\n4 Models", ACCENT_LIGHT),
        ("→", "", ACCENT),
        ("📈", "Evaluate\n& Compare", RED),
        ("→", "", ACCENT),
        ("🏆", "Best Model\n(Pickle)", GREEN),
    ]

    for i, (icon, label, color) in enumerate(steps):
        x = Inches(0.5 + i * 1.12)
        if icon == "→":
            add_text_box(slide, x, Inches(2.0), Inches(0.8), Inches(0.5),
                         "→", font_size=28, color=color, alignment=PP_ALIGN.CENTER)
        else:
            card = add_shape_box(slide, x, Inches(1.6), Inches(1.05), Inches(1.3))
            add_text_box(slide, x, Inches(1.65), Inches(1.05), Inches(0.5),
                         icon, font_size=24, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, x, Inches(2.1), Inches(1.05), Inches(0.7),
                         label, font_size=10, color=color, alignment=PP_ALIGN.CENTER)

    # Tech stack cards
    tech_categories = [
        ("🐍 Language", "Python 3.9+", GREEN),
        ("📊 Data", "Pandas, NumPy", ACCENT_LIGHT),
        ("🤖 ML", "Scikit-learn, XGBoost", YELLOW),
        ("⚖️ Balancing", "Imbalanced-learn\n(SMOTE)", RED),
        ("📈 Visualization", "Matplotlib, Seaborn\nPlotly", ACCENT_LIGHT),
        ("🌐 Dashboard", "Streamlit", GREEN),
    ]

    for i, (title, desc, color) in enumerate(tech_categories):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(3.4 + row * 1.7)
        card = add_shape_box(slide, x, y, Inches(3.8), Inches(1.4))
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.5),
                     title, font_size=16, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(0.7),
                     desc, font_size=14, color=DARK_TEXT)

    # ──────────────────────────────────────────
    #  SLIDE 6: Data Preprocessing
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "4. Data Preprocessing Pipeline", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    preprocessing_steps = [
        ("Step 1", "Load Dataset", "Read creditcard.csv using Pandas with error handling", "📥", GREEN),
        ("Step 2", "Inspect Data", "Check shape, dtypes, null counts, duplicates, class distribution", "🔍", ACCENT_LIGHT),
        ("Step 3", "Clean Data", "Remove ~1,081 duplicate rows; handle missing values with median", "🧹", YELLOW),
        ("Step 4", "Scale Features", "StandardScaler on Time & Amount (V1–V28 already PCA-scaled)", "📏", RED),
        ("Step 5", "Train/Test Split", "80/20 stratified split — maintains class proportions", "✂️", ACCENT_LIGHT),
        ("Step 6", "Apply SMOTE", "Oversample minority class on TRAINING data only (prevents leakage!)", "⚖️", GREEN),
    ]

    for i, (step, title, desc, icon, color) in enumerate(preprocessing_steps):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.5 + row * 1.85)

        card = add_shape_box(slide, x, y, Inches(5.8), Inches(1.6))
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.2), Inches(0.4),
                     f"{icon}  {step}: {title}", font_size=18, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.65), Inches(5.2), Inches(0.8),
                     desc, font_size=14, color=DARK_TEXT)

    # Important note
    card = add_shape_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8),
                         fill_color=RGBColor(0x2D, 0x14, 0x14), border_color=RED)
    add_text_box(slide, Inches(1.1), Inches(6.35), Inches(11.2), Inches(0.6),
                 "⚠️  CRITICAL: SMOTE is applied to training data ONLY — applying it before splitting causes data leakage!",
                 font_size=15, color=RED, bold=True)

    # ──────────────────────────────────────────
    #  SLIDE 7: SMOTE Explained
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "4b. SMOTE — Handling Class Imbalance", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5.5))

    # What is SMOTE card
    card = add_shape_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "What is SMOTE?", font_size=22, color=GREEN, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4.5),
                    [
                        "• Synthetic Minority Over-sampling Technique",
                        "• Creates SYNTHETIC fraud examples",
                        "• Interpolates between existing minority samples",
                        "• Balances the dataset to ~50/50 split",
                        "",
                        "Why not just duplicate fraud cases?",
                        "• Duplication causes overfitting",
                        "• SMOTE creates NEW realistic data points",
                        "• Better generalization on unseen data",
                    ], font_size=15)

    # Before/After card
    card = add_shape_box(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.4))
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.5),
                 "❌ Before SMOTE", font_size=20, color=RED, bold=True)
    add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(0.5),
                 "Legitimate: ~227,451 samples", font_size=16, color=DARK_TEXT)
    add_text_box(slide, Inches(7.1), Inches(2.7), Inches(5.2), Inches(0.5),
                 "Fraudulent:      ~394 samples", font_size=16, color=RED)
    add_text_box(slide, Inches(7.1), Inches(3.2), Inches(5.2), Inches(0.5),
                 "Ratio: 99.83% vs 0.17%", font_size=14, color=GRAY)

    card = add_shape_box(slide, Inches(6.8), Inches(4.2), Inches(5.7), Inches(2.8))
    add_text_box(slide, Inches(7.1), Inches(4.3), Inches(5), Inches(0.5),
                 "✅ After SMOTE", font_size=20, color=GREEN, bold=True)
    add_text_box(slide, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.5),
                 "Legitimate: ~227,451 samples", font_size=16, color=DARK_TEXT)
    add_text_box(slide, Inches(7.1), Inches(5.4), Inches(5.2), Inches(0.5),
                 "Fraudulent: ~227,451 samples (synthetic!)", font_size=16, color=GREEN)
    add_text_box(slide, Inches(7.1), Inches(5.9), Inches(5.2), Inches(0.5),
                 "Ratio: 50% vs 50% ✓", font_size=14, color=GREEN)
    add_text_box(slide, Inches(7.1), Inches(6.4), Inches(5.2), Inches(0.5),
                 "Applied to training data only!", font_size=13, color=YELLOW, bold=True)

    # ──────────────────────────────────────────
    #  SLIDE 8: EDA
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "5. Exploratory Data Analysis (EDA)", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    eda_items = [
        ("📊 Class Distribution", "Bar chart + pie chart showing 99.83% vs 0.17% imbalance", GREEN),
        ("💰 Amount Distribution", "Fraud amounts tend to be smaller (mean ~$122 vs ~$88)", YELLOW),
        ("⏰ Time Distribution", "Transaction frequency over the 48-hour recording period", ACCENT_LIGHT),
        ("🔗 Correlation Heatmap", "Feature correlations — V14, V12, V10 most correlated with fraud", RED),
        ("📈 Top Features", "Box plots of top 8 features showing class separation", ACCENT_LIGHT),
    ]

    for i, (title, desc, color) in enumerate(eda_items):
        y = Inches(1.5 + i * 1.1)
        card = add_shape_box(slide, Inches(0.8), y, Inches(11.7), Inches(0.95))
        add_text_box(slide, Inches(1.1), y + Inches(0.05), Inches(4), Inches(0.4),
                     title, font_size=18, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(11), Inches(0.4),
                     desc, font_size=14, color=DARK_TEXT)

    # Key insight box
    card = add_shape_box(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9),
                         fill_color=RGBColor(0x14, 0x2D, 0x1E), border_color=GREEN)
    add_text_box(slide, Inches(1.1), Inches(6.3), Inches(11.2), Inches(0.6),
                 "💡 Key Insight: Most fraud transactions are under $100, making them hard to detect by amount alone!",
                 font_size=16, color=GREEN, bold=True)

    # ──────────────────────────────────────────
    #  SLIDE 9: EDA Visualizations (images if available)
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "5b. EDA Visualizations", font_size=34, color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Try to embed actual EDA images
    assets_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
    images = [
        ("class_distribution.png", "Class Distribution"),
        ("amount_distribution.png", "Amount Distribution"),
    ]

    images_found = False
    for i, (img_file, title) in enumerate(images):
        img_path = os.path.join(assets_dir, img_file)
        if os.path.exists(img_path):
            x = Inches(0.5 + i * 6.3)
            slide.shapes.add_picture(img_path, x, Inches(1.5), Inches(6), Inches(3.5))
            add_text_box(slide, x, Inches(5.1), Inches(6), Inches(0.4),
                         title, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
            images_found = True

    if not images_found:
        card = add_shape_box(slide, Inches(1.5), Inches(2), Inches(10.3), Inches(4))
        add_text_box(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                     "📷 EDA plots will be embedded after running python3 train.py",
                     font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(2), Inches(4), Inches(9), Inches(0.6),
                     "Generated plots: class_distribution.png, amount_distribution.png,\n"
                     "time_distribution.png, correlation_heatmap.png, top_features.png",
                     font_size=15, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────
    #  SLIDE 10: ML Models
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "6. Machine Learning Models", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    models = [
        ("1️⃣ Logistic Regression",
         "Linear classifier — fast baseline",
         "• Simple, interpretable\n• Good for linearly separable data\n• Outputs probabilities\n• max_iter=1000, solver='lbfgs'",
         ACCENT_LIGHT),
        ("2️⃣ Decision Tree",
         "Tree-based — rule learning",
         "• Easy to interpret (if-else rules)\n• Handles non-linear patterns\n• Prone to overfitting\n• max_depth=10",
         GREEN),
        ("3️⃣ Random Forest",
         "Ensemble of 100 decision trees",
         "• Combines many trees (bagging)\n• Reduces overfitting\n• Robust and accurate\n• n_estimators=100, max_depth=15",
         YELLOW),
        ("4️⃣ XGBoost",
         "Gradient boosted trees",
         "• Trees learn from previous errors\n• State-of-the-art performance\n• Handles imbalance well\n• n_estimators=100, lr=0.1",
         RED),
    ]

    for i, (title, subtitle, details, color) in enumerate(models):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.5 + row * 2.85)

        card = add_shape_box(slide, x, y, Inches(5.8), Inches(2.6))
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.2), Inches(0.5),
                     title, font_size=20, color=color, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.2), Inches(0.4),
                     subtitle, font_size=13, color=GRAY)
        add_text_box(slide, x + Inches(0.2), y + Inches(1.0), Inches(5.2), Inches(1.4),
                     details, font_size=13, color=DARK_TEXT)

    # ──────────────────────────────────────────
    #  SLIDE 11: Evaluation Metrics
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "7. Evaluation Metrics Explained", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    # Why not accuracy card
    card = add_shape_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
                         fill_color=RGBColor(0x2D, 0x14, 0x14), border_color=RED)
    add_text_box(slide, Inches(1.1), Inches(1.55), Inches(11.2), Inches(0.5),
                 "⚠️ Why NOT Accuracy?", font_size=20, color=RED, bold=True)
    add_text_box(slide, Inches(1.1), Inches(2.05), Inches(11.2), Inches(0.5),
                 "A model predicting ALL transactions as legitimate achieves 99.83% accuracy — but catches ZERO fraud!",
                 font_size=15, color=DARK_TEXT)

    # Metrics grid
    metrics = [
        ("Accuracy", "Overall correct predictions\nTP + TN / Total", "Good for balanced datasets only", GRAY),
        ("Precision", "Of predicted frauds, how many\nare actual frauds?", "Minimizes false alarms", ACCENT_LIGHT),
        ("Recall", "Of actual frauds, how many\nwere detected?", "Catches the most frauds", GREEN),
        ("F1 Score ⭐", "Harmonic mean of Precision\nand Recall", "PRIMARY METRIC — balances both!", YELLOW),
        ("ROC-AUC", "Model's ability to distinguish\nbetween classes", "Ranking & threshold metric", ACCENT_LIGHT),
        ("Confusion Matrix", "Table of TP, TN, FP, FN\ncounts", "Visual error breakdown", RED),
    ]

    for i, (name, desc, note, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.1)
        y = Inches(3.0 + row * 2.1)

        card = add_shape_box(slide, x, y, Inches(3.8), Inches(1.85))
        add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(3.5), Inches(0.4),
                     name, font_size=17, color=color, bold=True)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(0.7),
                     desc, font_size=13, color=DARK_TEXT)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.35), Inches(3.5), Inches(0.4),
                     f"→ {note}", font_size=11, color=GRAY)

    # ──────────────────────────────────────────
    #  SLIDE 12: Results
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "8. Results & Model Comparison", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    # Results table header
    table_data = [
        ("Model", "Accuracy", "Precision", "Recall", "F1 Score", "ROC-AUC"),
        ("Logistic Regression", "~0.974", "~0.06", "~0.92", "~0.11", "~0.97"),
        ("Decision Tree", "~0.997", "~0.72", "~0.76", "~0.74", "~0.88"),
        ("Random Forest", "~0.999", "~0.93", "~0.80", "~0.86", "~0.96"),
        ("XGBoost", "~0.999", "~0.88", "~0.82", "~0.85", "~0.97"),
    ]

    # Draw table
    col_widths = [Inches(2.8), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)]
    table_left = Inches(0.8)
    table_top = Inches(1.5)
    row_height = Inches(0.55)

    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            x = table_left + sum(w for w in col_widths[:col_idx])
            y = table_top + row_idx * row_height

            if row_idx == 0:
                # Header row
                card = add_shape_box(slide, x, y, col_widths[col_idx] - Pt(2), row_height - Pt(2),
                                     fill_color=ACCENT, border_color=ACCENT)
                add_text_box(slide, x + Pt(4), y + Pt(2), col_widths[col_idx] - Pt(8), row_height - Pt(4),
                             cell_text, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
            else:
                bg = BG_CARD if row_idx % 2 == 1 else RGBColor(0x1A, 0x27, 0x48)
                border = ACCENT if row_idx == 4 else RGBColor(0x2D, 0x35, 0x61)
                card = add_shape_box(slide, x, y, col_widths[col_idx] - Pt(2), row_height - Pt(2),
                                     fill_color=bg, border_color=border)
                text_color = GREEN if (row_idx == 3 and col_idx == 4) else DARK_TEXT
                add_text_box(slide, x + Pt(4), y + Pt(2), col_widths[col_idx] - Pt(8), row_height - Pt(4),
                             cell_text, font_size=13, color=text_color,
                             bold=(row_idx == 3 and col_idx >= 1),
                             alignment=PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT)

    # Best model highlight
    card = add_shape_box(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(1.2),
                         fill_color=RGBColor(0x14, 0x2D, 0x1E), border_color=GREEN)
    add_text_box(slide, Inches(1.1), Inches(4.55), Inches(11), Inches(0.5),
                 "🏆 Best Model: Random Forest / XGBoost", font_size=22, color=GREEN, bold=True)
    add_text_box(slide, Inches(1.1), Inches(5.1), Inches(11), Inches(0.5),
                 "Selected based on highest F1 Score — the best metric for imbalanced classification problems",
                 font_size=15, color=DARK_TEXT)

    # Analysis points
    card = add_shape_box(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.2))
    add_text_box(slide, Inches(1.1), Inches(6.05), Inches(11), Inches(0.4),
                 "📊 Key Observations", font_size=18, color=ACCENT_LIGHT, bold=True)
    add_text_box(slide, Inches(1.1), Inches(6.45), Inches(11), Inches(0.6),
                 "• LR: High recall but very low precision (too many false alarms)  •  Ensemble methods (RF, XGB) far outperform single models  •  Accuracy alone is misleading",
                 font_size=13, color=DARK_TEXT)

    # ──────────────────────────────────────────
    #  SLIDE 13: Streamlit Dashboard
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "9. Streamlit Web Dashboard", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    # Dashboard features
    pages = [
        ("🏠 Home Page", [
            "• Project overview with summary metric cards",
            "• Dataset statistics (284K transactions, 0.172% fraud)",
            "• Best model highlight card",
            "• All 5 EDA visualizations in interactive tabs",
            "• Key insights from the data",
        ], GREEN),
        ("🔍 Prediction Page", [
            "• Input form for all 30 features (V1–V28, Time, Amount)",
            "• Quick-fill buttons for demo (legitimate & fraud samples)",
            "• Real-time fraud prediction with confidence percentage",
            "• Animated result cards (green = safe, red = fraud)",
            "• Gauge chart for fraud probability visualization",
        ], ACCENT_LIGHT),
        ("📊 Model Performance", [
            "• Side-by-side comparison table of all 4 models",
            "• Interactive radar chart for best model metrics",
            "• Grouped bar chart comparison",
            "• Confusion matrices and ROC curves",
        ], YELLOW),
        ("ℹ️ About Page", [
            "• Tech stack information",
            "• ML pipeline explanation",
            "• Dataset documentation",
            "• Project credits",
        ], RED),
    ]

    for i, (title, features, color) in enumerate(pages):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.5 + row * 2.85)

        card = add_shape_box(slide, x, y, Inches(5.8), Inches(2.6))
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.2), Inches(0.5),
                     title, font_size=20, color=color, bold=True)
        add_bullet_list(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.2), Inches(1.8),
                        features, font_size=12, color=DARK_TEXT)

    # ──────────────────────────────────────────
    #  SLIDE 14: Conclusion & Future Work
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
                 "10. Conclusion & Future Work", font_size=34,
                 color=ACCENT_LIGHT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(5))

    # Conclusion card
    card = add_shape_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
    add_text_box(slide, Inches(1.1), Inches(1.6), Inches(5), Inches(0.5),
                 "✅ What We Achieved", font_size=22, color=GREEN, bold=True)
    add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5), Inches(4.5),
                    [
                        "• Complete ML pipeline for fraud detection",
                        "• Handled extreme class imbalance (0.172%)",
                        "    using SMOTE oversampling",
                        "• Trained & compared 4 ML algorithms",
                        "• Achieved F1 Score > 0.85 with ensemble models",
                        "• Deployed production-style Streamlit dashboard",
                        "• Real-time prediction with confidence scores",
                        "• Modular, well-documented, reusable code",
                        "• University-ready documentation & report",
                    ], font_size=14)

    # Future work card
    card = add_shape_box(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.5))
    add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5), Inches(0.5),
                 "🚀 Future Improvements", font_size=22, color=YELLOW, bold=True)
    add_bullet_list(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.5),
                    [
                        "• Deep Learning: LSTM / Autoencoders",
                        "    for sequential pattern detection",
                        "• Real-time streaming: Apache Kafka",
                        "    integration for live processing",
                        "• SHAP values for model explainability",
                        "    (transparent, trustworthy predictions)",
                        "• Hyperparameter tuning with Optuna",
                        "• REST API with FastAPI for production",
                        "• Concept drift monitoring & auto-retrain",
                    ], font_size=14)

    # ──────────────────────────────────────────
    #  SLIDE 15: Thank You
    # ──────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.2),
                 "Thank You!", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_accent_line(slide, Inches(5), Inches(3.3), Inches(3.5))

    add_text_box(slide, Inches(2), Inches(3.8), Inches(9.5), Inches(0.8),
                 "Credit Card Fraud Detection using Machine Learning",
                 font_size=22, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(4.8), Inches(9.5), Inches(0.6),
                 "Questions & Discussion",
                 font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Pt(6), prs.slide_width, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    return prs


# ──────────────────────────────────────────────
#  Main
# ──────────────────────────────────────────────
if __name__ == "__main__":
    print("🎨 Generating presentation...")
    prs = create_presentation()

    output_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "Credit_Card_Fraud_Detection_Presentation.pptx"
    )
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
